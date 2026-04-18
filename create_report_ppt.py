#!/usr/bin/env python3
"""生成项目1&2汇报PPT，并重绘项目二 SVG 图示。"""

import os
from pathlib import Path

os.environ.setdefault("MPLCONFIGDIR", "/tmp/mplconfig_codex")

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Polygon
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE = "/home/siton02/md0/crf/cl_zdjs"
CONTENT_W = 11.7
CONTENT_H = 5.45
ASSET_DIR = Path(BASE) / "project2_centerpoint" / "generated_assets"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── PPT Color scheme ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_GREEN = RGBColor(0x00, 0xC9, 0x8D)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
CARD_COLOR = RGBColor(0x22, 0x2B, 0x45)

# ── Matplotlib Color scheme ──
BG_DARK_HEX = "#1A1A2E"
CARD_HEX = "#222B45"
PANEL_HEX = "#202943"
PANEL_LINE_HEX = "#35557B"
ACCENT_BLUE_HEX = "#0096D6"
ACCENT_GREEN_HEX = "#00C98D"
ACCENT_ORANGE_HEX = "#FF8C00"
WHITE_HEX = "#FFFFFF"
LIGHT_GRAY_HEX = "#CCD4E0"
MUTED_HEX = "#9FB0C7"


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="微软雅黑"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = color
    paragraph.font.bold = bold
    paragraph.font.name = font_name
    paragraph.alignment = align
    return box


def add_accent_bar(slide, left, top, width=0.08, height=0.6, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bottom_bar(slide, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2),
                                   Inches(13.333), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_card(slide, left, top, width, height, color=CARD_COLOR):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=LIGHT_GRAY):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.font.name = "微软雅黑"
        paragraph.space_after = Pt(8)
    return box


def create_asset_figure():
    plt.rcParams["font.sans-serif"] = ["Noto Sans CJK SC", "WenQuanYi Zen Hei", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False
    fig, ax = plt.subplots(figsize=(CONTENT_W, CONTENT_H), dpi=220)
    fig.patch.set_facecolor(BG_DARK_HEX)
    ax.set_facecolor(BG_DARK_HEX)
    ax.set_xlim(0, CONTENT_W)
    ax.set_ylim(CONTENT_H, 0)
    ax.axis("off")
    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)
    return fig, ax


def add_mpl_card(ax, x, y, width, height, face_color=CARD_HEX,
                 edge_color=None, linewidth=1.6, radius=0.18):
    card = FancyBboxPatch(
        (x, y),
        width,
        height,
        boxstyle=f"round,pad=0.02,rounding_size={radius}",
        linewidth=linewidth,
        facecolor=face_color,
        edgecolor=edge_color or face_color,
    )
    ax.add_patch(card)
    return card


def add_mpl_text(ax, x, y, text, size=14, color=WHITE_HEX, weight="normal",
                 ha="left", va="center", alpha=1.0, linespacing=1.3):
    return ax.text(
        x,
        y,
        text,
        fontsize=size,
        color=color,
        fontweight=weight,
        ha=ha,
        va=va,
        alpha=alpha,
        linespacing=linespacing,
    )


def add_chevron(ax, x, y, width=0.26, height=0.22, color=ACCENT_GREEN_HEX):
    points = [
        (x, y - height),
        (x + width * 0.55, y - height),
        (x + width, y),
        (x + width * 0.55, y + height),
        (x, y + height),
        (x + width * 0.35, y),
    ]
    ax.add_patch(Polygon(points, closed=True, facecolor=color, edgecolor=color))


def add_image_to_axis(ax, img_path, x, y, width, height):
    image = plt.imread(str(img_path))
    ax.imshow(image, extent=[x, x + width, y + height, y], aspect="auto", zorder=2)


def save_asset_figure(fig, stem):
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    svg_path = ASSET_DIR / f"{stem}.svg"
    png_path = ASSET_DIR / f"{stem}.png"
    fig.savefig(svg_path, facecolor=fig.get_facecolor())
    fig.savefig(png_path, facecolor=fig.get_facecolor())
    plt.close(fig)
    return str(png_path), str(svg_path)


def generate_centerpoint_structure_assets():
    fig, ax = create_asset_figure()

    add_mpl_card(ax, 0.08, 0.08, 11.54, 5.29, PANEL_HEX, PANEL_LINE_HEX, 1.6, 0.22)
    add_mpl_text(ax, 0.36, 0.48, "CenterPoint 结构流程", 20, WHITE_HEX, "bold")
    add_mpl_text(ax, 0.36, 0.84, "SVG 矢量重绘版：结构分层展示，避免文字重叠和连线杂乱", 10.0, MUTED_HEX)

    stage_pills = [
        (0.46, 1.12, 1.55, "输入层"),
        (2.63, 1.12, 4.12, "编码层"),
        (6.97, 1.12, 1.95, "预测层"),
        (9.14, 1.12, 1.95, "输出层"),
    ]
    for x, y, width, label in stage_pills:
        add_mpl_card(ax, x, y, width, 0.38, "#1D3554", ACCENT_GREEN_HEX, 1.0, 0.12)
        add_mpl_text(ax, x + width / 2, y + 0.19, label, 11.2, ACCENT_GREEN_HEX, "bold", "center")

    node_specs = [
        ("01", "点云输入", "LiDAR sweep\n多帧同步与融合", 0.46, ACCENT_BLUE_HEX),
        ("02", "体素化编码", "dynamic voxel\n保留稀疏结构", 2.63, "#35A7FF"),
        ("03", "特征提取", "VFE + sparse\nbackbone / neck", 4.80, ACCENT_GREEN_HEX),
        ("04", "检测头", "center heatmap\nsize / yaw / velocity", 6.97, "#55D977"),
        ("05", "3D框输出", "decode + NMS\n位置 / 尺寸 / 朝向", 9.14, ACCENT_ORANGE_HEX),
    ]
    node_y = 1.66
    node_w = 1.95
    node_h = 1.50
    for index, title, desc, x, accent in node_specs:
        add_mpl_card(ax, x, node_y, node_w, node_h, "#24304D", accent, 1.8, 0.18)
        add_mpl_text(ax, x + 0.18, node_y + 0.28, index, 10.5, accent, "bold")
        add_mpl_text(ax, x + 0.18, node_y + 0.66, title, 14.6, WHITE_HEX, "bold")
        add_mpl_text(ax, x + 0.18, node_y + 1.10, desc, 10.4, LIGHT_GRAY_HEX, va="center")

    for x in (2.26, 4.43, 6.60, 8.77):
        add_chevron(ax, x, 2.41, 0.24, 0.20, ACCENT_GREEN_HEX)

    summary_cards = [
        (0.58, "输入数据", "多帧 LiDAR 点云\n时间同步与预处理"),
        (4.02, "检测核心", "中心热力图监督\n多任务属性回归"),
        (7.46, "部署价值", "与 2D 感知互补\n输出可直接服务定位与决策"),
    ]
    for x, title, desc in summary_cards:
        add_mpl_card(ax, x, 3.54, 3.10, 1.06, "#1F2942", PANEL_LINE_HEX, 1.2, 0.16)
        add_mpl_text(ax, x + 0.16, 3.84, title, 13.0, ACCENT_GREEN_HEX, "bold")
        add_mpl_text(ax, x + 0.16, 4.22, desc, 10.0, LIGHT_GRAY_HEX, va="center")

    add_mpl_card(ax, 0.58, 4.76, 10.04, 0.46, "#1A2236", "#1A2236", 0.8, 0.12)
    add_mpl_text(
        ax,
        0.78,
        4.99,
        "输出字段：位置 (x, y, z) / 尺寸 (w, l, h) / 朝向 yaw / 速度 (vx, vy)",
        9.8,
        MUTED_HEX,
    )

    return save_asset_figure(fig, "centerpoint_structure")


def generate_centerpoint_validation_assets():
    fig, ax = create_asset_figure()
    detect_path = Path(BASE) / "project2_centerpoint" / "CenterPoint" / "detection_result.jpg"
    gt_path = Path(BASE) / "project2_centerpoint" / "CenterPoint" / "gt_result.jpg"

    add_mpl_card(ax, 0.08, 0.08, 11.54, 5.29, PANEL_HEX, PANEL_LINE_HEX, 1.6, 0.22)
    add_mpl_text(ax, 0.36, 0.48, "检测结果对比", 20, WHITE_HEX, "bold")
    add_mpl_text(ax, 0.36, 0.84, "说明文字统一移到图外，保留图像主体，避免图内信息拥挤", 10.0, MUTED_HEX)

    image_cards = [
        (0.38, "检测场景图", detect_path, "原始前视场景，便于观察道路与目标分布"),
        (6.02, "标注结果对比", gt_path, "绿色框为标注结果，空间定位关系更直观"),
    ]
    for x, title, img_path, caption in image_cards:
        add_mpl_card(ax, x, 1.08, 5.30, 2.80, "#24304D", ACCENT_GREEN_HEX, 1.3, 0.16)
        add_mpl_text(ax, x + 0.16, 1.34, title, 13.6, ACCENT_GREEN_HEX, "bold")
        add_image_to_axis(ax, img_path, x + 0.16, 1.56, 4.98, 1.82)
        add_mpl_text(ax, x + 0.16, 3.62, caption, 9.6, LIGHT_GRAY_HEX)

    summary_cards = [
        ("验证数据", "nuScenes mini-val\n多帧 LiDAR 与前视图像对齐"),
        ("检测能力", "车辆、行人等目标的\n3D 空间位置与尺度估计"),
        ("当前状态", "代码调试与权重准备完成\n下一步进行车载部署适配"),
    ]
    summary_x = [0.50, 4.08, 7.66]
    for x, (title, desc) in zip(summary_x, summary_cards):
        add_mpl_card(ax, x, 4.04, 3.06, 1.06, "#1F2942", PANEL_LINE_HEX, 1.2, 0.16)
        add_mpl_text(ax, x + 0.16, 4.31, title, 12.6, ACCENT_GREEN_HEX, "bold")
        add_mpl_text(ax, x + 0.16, 4.74, desc, 9.8, LIGHT_GRAY_HEX, va="center")

    return save_asset_figure(fig, "centerpoint_validation")


centerpoint_structure_png, centerpoint_structure_svg = generate_centerpoint_structure_assets()
centerpoint_validation_png, centerpoint_validation_svg = generate_centerpoint_validation_assets()


# ════════════════════════════════════════════════════════════════
# Slide 1: Title
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               Inches(13.333), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

add_text_box(slide, 1.5, 1.8, 10, 1.2, "智能驾驶感知算法", 44, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 2.9, 10, 1.0, "项目进展汇报", 36, WHITE, True, PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.0),
                               Inches(3.333), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

add_text_box(slide, 1.5, 4.3, 10, 0.6,
             "项目一：基于YOLO的2D障碍物检测  |  项目二：基于CenterPoint的3D目标检测",
             18, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 5.5, 10, 0.5, "汇报日期：2026年3月17日", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_bottom_bar(slide, ACCENT_BLUE)


# ════════════════════════════════════════════════════════════════
# Slide 2: Outline
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_accent_bar(slide, 0.8, 0.6, 0.08, 0.5, ACCENT_BLUE)
add_text_box(slide, 1.0, 0.55, 5, 0.6, "汇报提纲", 30, WHITE, True)

items_data = [
    ("01", "项目一：基于YOLO的2D障碍物检测", "算法原理与实验验证", ACCENT_BLUE),
    ("02", "项目二：基于CenterPoint的3D目标检测", "结构图重绘与结果验证", ACCENT_GREEN),
    ("03", "实验平台搭建", "智能小车组装完成", ACCENT_ORANGE),
    ("04", "下一步计划", "算法部署与车载验证", RGBColor(0xE0, 0x56, 0x7E)),
]

for i, (num, title, desc, color) in enumerate(items_data):
    y = 1.8 + i * 1.3
    add_card(slide, 1.5, y, 10, 1.0)
    add_text_box(slide, 1.8, y + 0.1, 0.8, 0.8, num, 28, color, True)
    add_text_box(slide, 2.8, y + 0.08, 7, 0.5, title, 22, WHITE, True)
    add_text_box(slide, 2.8, y + 0.55, 7, 0.4, desc, 14, LIGHT_GRAY)

add_bottom_bar(slide, ACCENT_BLUE)


# ════════════════════════════════════════════════════════════════
# Slide 3: Project 1 - Algorithm Overview
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_accent_bar(slide, 0.8, 0.6, 0.08, 0.5, ACCENT_BLUE)
add_text_box(slide, 1.0, 0.55, 8, 0.6, "项目一：基于YOLO的2D障碍物检测", 28, WHITE, True)

add_card(slide, 0.8, 1.5, 5.8, 5.2)
add_text_box(slide, 1.1, 1.7, 5, 0.5, "算法概述", 20, ACCENT_BLUE, True)
add_bullet_text(slide, 1.1, 2.3, 5.2, 4.0, [
    "  采用 YOLOv5 与 YOLO11 系列模型",
    "  针对自动驾驶场景进行障碍物检测",
    "  支持多类别目标：车辆、行人、卡车、",
    "  自行车、摩托车、交通灯等",
    "  在 nuScenes 数据集上完成验证",
    "  模型对比：n / s / m / l 四种规格",
    "  兼顾检测精度与推理速度",
], 15, LIGHT_GRAY)

add_card(slide, 7.0, 1.5, 5.5, 5.2)
add_text_box(slide, 7.3, 1.7, 5, 0.5, "核心指标", 20, ACCENT_BLUE, True)

metrics = [
    ("检测帧数", "81 帧"),
    ("总检测数", "1,107 个目标"),
    ("平均推理", "56.0 ms / 帧"),
    ("实时帧率", "17.8 FPS"),
    ("动态目标", "车辆833 / 行人161 / 卡车30"),
    ("静态目标", "交通灯 31"),
]
for i, (label, value) in enumerate(metrics):
    y = 2.4 + i * 0.7
    add_text_box(slide, 7.3, y, 2.2, 0.4, label, 14, LIGHT_GRAY, True)
    add_text_box(slide, 9.3, y, 3.0, 0.4, value, 14, ACCENT_GREEN)

add_bottom_bar(slide, ACCENT_BLUE)


# ════════════════════════════════════════════════════════════════
# Slide 4: Project 1 - Detection Results
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_accent_bar(slide, 0.8, 0.6, 0.08, 0.5, ACCENT_BLUE)
add_text_box(slide, 1.0, 0.55, 8, 0.6, "项目一：YOLO 检测效果展示", 28, WHITE, True)

img_path = os.path.join(BASE, "project1_yolo_obstacle/results/detect_front_00.jpg")
slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.5), Inches(7.5), Inches(4.5))
add_text_box(slide, 0.8, 6.1, 7.5, 0.5,
             "YOLO11m 检测结果：识别车辆、行人、交通灯等多类别目标", 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_card(slide, 8.8, 1.5, 3.8, 5.2)
add_text_box(slide, 9.1, 1.7, 3.4, 0.5, "检测能力", 18, ACCENT_BLUE, True)
add_bullet_text(slide, 9.1, 2.4, 3.4, 4.0, [
    "  动态目标检测",
    "  - 车辆 (car)",
    "  - 行人 (person)",
    "  - 卡车 (truck)",
    "  - 自行车 (bicycle)",
    "  - 摩托车 (motorcycle)",
    "",
    "  静态目标检测",
    "  - 交通灯 (traffic light)",
    "",
    "  多置信度阈值自适应",
    "  动态/静态分类标注",
], 13, LIGHT_GRAY)

add_bottom_bar(slide, ACCENT_BLUE)


# ════════════════════════════════════════════════════════════════
# Slide 5: Project 1 - Model Comparison
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_accent_bar(slide, 0.8, 0.6, 0.08, 0.5, ACCENT_BLUE)
add_text_box(slide, 1.0, 0.55, 8, 0.6, "项目一：模型性能对比分析", 28, WHITE, True)

img_path = os.path.join(BASE, "project1_yolo_obstacle/results/model_comparison.jpg")
slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.4), Inches(8.0), Inches(2.8))
add_text_box(slide, 0.8, 4.3, 8, 0.4, "YOLO11 不同规格模型检测效果对比 (n/s/m/l)", 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)

img_path2 = os.path.join(BASE, "project1_yolo_obstacle/results/speed_accuracy_tradeoff.png")
slide.shapes.add_picture(img_path2, Inches(0.8), Inches(4.7), Inches(5.5), Inches(2.5))

add_card(slide, 9.2, 1.4, 3.5, 5.8)
add_text_box(slide, 9.4, 1.6, 3.0, 0.4, "模型对比", 18, ACCENT_BLUE, True)

models_info = [
    ("YOLO11n", "2.6M", "22ms", "45 FPS"),
    ("YOLO11s", "9.5M", "23ms", "43 FPS"),
    ("YOLO11m", "20.1M", "25ms", "40 FPS"),
    ("YOLO11l", "25.4M", "35ms", "28 FPS"),
]
add_text_box(slide, 9.4, 2.2, 1.2, 0.35, "模型", 11, ACCENT_GREEN, True)
add_text_box(slide, 10.3, 2.2, 0.8, 0.35, "参数量", 11, ACCENT_GREEN, True)
add_text_box(slide, 11.0, 2.2, 0.7, 0.35, "延迟", 11, ACCENT_GREEN, True)
add_text_box(slide, 11.6, 2.2, 1.0, 0.35, "帧率", 11, ACCENT_GREEN, True)

for i, (name, params, latency, fps) in enumerate(models_info):
    y = 2.6 + i * 0.45
    add_text_box(slide, 9.4, y, 1.2, 0.35, name, 12, WHITE)
    add_text_box(slide, 10.3, y, 0.8, 0.35, params, 12, LIGHT_GRAY)
    add_text_box(slide, 11.0, y, 0.7, 0.35, latency, 12, LIGHT_GRAY)
    add_text_box(slide, 11.6, y, 1.0, 0.35, fps, 12, ACCENT_GREEN)

add_text_box(slide, 9.4, 4.5, 3.0, 1.5,
             "结论：YOLO11s 在精度与速度之间\n取得最佳平衡，推荐用于车载部署。\n\nYOLO11n 适合资源受限的嵌入式\n平台（如 Jetson Nano）。",
             12, LIGHT_GRAY)
add_bottom_bar(slide, ACCENT_BLUE)


# ════════════════════════════════════════════════════════════════
# Slide 6: Project 2 - SVG Structure Diagram
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_accent_bar(slide, 0.8, 0.6, 0.08, 0.5, ACCENT_GREEN)
add_text_box(slide, 1.0, 0.55, 8, 0.6, "项目二：CenterPoint 结构图", 28, WHITE, True)
add_text_box(slide, 1.0, 1.02, 9.8, 0.35, "已改为 SVG 矢量重绘，去掉杂乱连线并重新整理结构层次", 12, LIGHT_GRAY)
slide.shapes.add_picture(centerpoint_structure_png, Inches(0.8), Inches(1.38), Inches(CONTENT_W), Inches(CONTENT_H))
add_bottom_bar(slide, ACCENT_GREEN)


# ════════════════════════════════════════════════════════════════
# Slide 7: Project 2 - Clean Comparison Figure
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_accent_bar(slide, 0.8, 0.6, 0.08, 0.5, ACCENT_GREEN)
add_text_box(slide, 1.0, 0.55, 8, 0.6, "项目二：CenterPoint 检测效果展示", 28, WHITE, True)
add_text_box(slide, 1.0, 1.02, 9.8, 0.35, "结果图说明已移到图外，避免图内文字挤压主体内容", 12, LIGHT_GRAY)
slide.shapes.add_picture(centerpoint_validation_png, Inches(0.8), Inches(1.38), Inches(CONTENT_W), Inches(CONTENT_H))
add_bottom_bar(slide, ACCENT_GREEN)


# ════════════════════════════════════════════════════════════════
# Slide 8: Vehicle Assembly - JetBot
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_accent_bar(slide, 0.8, 0.6, 0.08, 0.5, ACCENT_ORANGE)
add_text_box(slide, 1.0, 0.55, 8, 0.6, "实验平台：智能小车组装完成", 28, WHITE, True)

img_car = os.path.join(BASE, "bdf02831db87ad16aa21b8b6e0d0a091.jpg")
slide.shapes.add_picture(img_car, Inches(1.5), Inches(1.4), Inches(5.5), Inches(4.5))
add_text_box(slide, 1.5, 6.0, 5.5, 0.4, "基于 NVIDIA Jetson Nano 的智能小车平台", 13, LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_card(slide, 7.5, 1.4, 5.0, 5.2)
add_text_box(slide, 7.8, 1.6, 4.5, 0.5, "平台配置", 20, ACCENT_ORANGE, True)
add_bullet_text(slide, 7.8, 2.3, 4.5, 4.0, [
    "  计算平台",
    "  - NVIDIA Jetson Nano",
    "  - 128核 Maxwell GPU",
    "  - 4GB 内存",
    "",
    "  传感器",
    "  - CSI 摄像头（用于2D检测）",
    "  - WiFi 天线（远程通信）",
    "",
    "  执行机构",
    "  - 双轮差速驱动底盘",
    "  - PWM 电机控制",
    "",
    "  当前状态：组装完成，待算法移植",
], 14, LIGHT_GRAY)

add_card(slide, 7.8, 6.0, 4.2, 0.5, ACCENT_GREEN)
add_text_box(slide, 7.8, 6.05, 4.2, 0.4, "硬件平台已就绪", 16, WHITE, True, PP_ALIGN.CENTER)
add_bottom_bar(slide, ACCENT_ORANGE)


# ════════════════════════════════════════════════════════════════
# Slide 9: Summary & Next Steps
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

add_accent_bar(slide, 0.8, 0.6, 0.08, 0.5, ACCENT_BLUE)
add_text_box(slide, 1.0, 0.55, 8, 0.6, "总结与下一步计划", 28, WHITE, True)

add_card(slide, 0.8, 1.5, 5.8, 5.2)
add_text_box(slide, 1.1, 1.7, 5, 0.5, "当前进展", 20, ACCENT_BLUE, True)

progress_items = [
    ("项目一 - YOLO 2D检测", "算法验证完成", ACCENT_GREEN),
    ("项目二 - CenterPoint 3D检测", "结构图与结果页优化", ACCENT_GREEN),
    ("智能小车平台", "组装完成", ACCENT_GREEN),
    ("算法车载部署", "待进行", ACCENT_ORANGE),
]
for i, (item, status, color) in enumerate(progress_items):
    y = 2.4 + i * 0.9
    add_text_box(slide, 1.3, y, 3.5, 0.35, item, 14, WHITE)
    add_card(slide, 4.8, y, 1.5, 0.35, color)
    add_text_box(slide, 4.8, y, 1.5, 0.35, status, 12, WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, 1.3, 6.0, 5, 0.5, "整体进度：感知算法与汇报材料均已整理，即将进入部署阶段", 13, ACCENT_GREEN, True)

add_card(slide, 7.0, 1.5, 5.5, 5.2)
add_text_box(slide, 7.3, 1.7, 5, 0.5, "下一步计划", 20, ACCENT_ORANGE, True)
add_bullet_text(slide, 7.3, 2.4, 5.0, 4.5, [
    "  1. YOLO模型移植到 Jetson Nano",
    "     - TensorRT 模型优化加速",
    "     - 摄像头实时推理测试",
    "",
    "  2. CenterPoint 算法适配",
    "     - 模型轻量化与推理优化",
    "     - 传感器数据接口对接",
    "",
    "  3. 系统集成与测试",
    "     - 感知-决策-控制闭环",
    "     - 实车场景验证测试",
    "",
    "  4. 性能优化与报告",
    "     - 实时性能调优",
    "     - 实验数据收集与分析",
], 13, LIGHT_GRAY)

add_bottom_bar(slide, ACCENT_BLUE)


# ════════════════════════════════════════════════════════════════
# Slide 10: Thank you
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               Inches(13.333), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

add_text_box(slide, 1.5, 2.5, 10, 1.0, "谢谢！", 48, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.8, 10, 0.6, "欢迎提问与讨论", 24, LIGHT_GRAY, False, PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.8),
                               Inches(3.333), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()
add_bottom_bar(slide, ACCENT_BLUE)


output_path = os.path.join(BASE, "项目汇报_智能驾驶感知算法.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"SVG assets: {centerpoint_structure_svg}, {centerpoint_validation_svg}")
